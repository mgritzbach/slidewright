#!/usr/bin/env python3
"""Create a deliberately invalid all-raster deck from per-slide PNGs."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("preview_dir", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()
    images = sorted(args.preview_dir.glob("slide-*.png"))
    if not images:
        parser.error("preview_dir has no slide-*.png images")
    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    while deck.slides:
        relationship_id = deck.slides._sldIdLst[-1].rId  # pylint: disable=protected-access
        deck.part.drop_rel(relationship_id)
        del deck.slides._sldIdLst[-1]  # pylint: disable=protected-access
    blank = deck.slide_layouts[6]
    for image in images:
        slide = deck.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0, width=deck.slide_width, height=deck.slide_height)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    deck.save(args.output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
